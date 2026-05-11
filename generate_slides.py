from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

FIGS = "artifacts/figures"

BLUE_DARK = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_MID = RGBColor(0x2C, 0x5F, 0x8A)
BLUE_LIGHT = RGBColor(0x4A, 0x90, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
ACCENT = RGBColor(0xE8, 0x8D, 0x2A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=18, color=BLACK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return txBox

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)

def add_title_bar(slide, text):
    add_rect(slide, 0, 0, W, Inches(1.1), BLUE_DARK)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.8), text,
                 font_size=32, color=WHITE, bold=True)
    add_rect(slide, 0, Inches(1.1), W, Inches(0.06), ACCENT)

def add_footer(slide):
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), BLUE_DARK)
    add_text_box(slide, Inches(0.5), Inches(7.12), Inches(12), Inches(0.35),
                 "House Price Prediction  ·  AI Foundations URV",
                 font_size=11, color=WHITE)

def add_info_box(slide, left, top, width, height, text, bg=LIGHT_BG, text_color=BLACK, bold_part=""):
    add_rect(slide, left, top, width, height, bg)
    if bold_part:
        add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), height - Inches(0.16),
                     text, font_size=14, color=text_color)
    else:
        add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), height - Inches(0.16),
                     text, font_size=14, color=text_color)

def add_card(slide, left, top, width, height, title, body, bar_color=BLUE_LIGHT):
    add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, Inches(0.06), height, bar_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
                 title, font_size=17, color=bar_color, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                 body, font_size=14, color=GRAY)

# ═══════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_rect(slide, 0, Inches(2.5), W, Inches(2.8), RGBColor(0x24, 0x4D, 0x75))
add_rect(slide, 0, Inches(3.6), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.7),
             "🏠  House Price Prediction", font_size=48, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Mi primer proyecto de Machine Learning",
             font_size=22, color=RGBColor(0xCC, 0xD5, 0xE0), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "Artificial Intelligence Foundations — Fundació URV",
             font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.4),
             "Adrià  ·  Raul  |  Mayo 2026",
             font_size=16, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 2 — QUÉ HICIMOS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "🎯  ¿Qué queríamos conseguir?")
add_footer(slide)

add_card(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(2.6),
         "El problema",
         "Predecir el precio de una casa (SalePrice)\n"
         "a partir de sus características\n"
         "(tamaño, calidad, ubicación...)", BLUE_MID)

add_card(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(2.6),
         "El dataset",
         "Kaggle House Prices\n"
         "1.460 casas · 80 variables\n"
         "Datos reales de Ames, Iowa", BLUE_LIGHT)

add_card(slide, Inches(0.5), Inches(4.3), Inches(6), Inches(2.5),
         "El enfoque",
         "Pipeline completo de ML:\n"
         "EDA → Preprocesamiento → Modelo → App", BLUE_MID)

add_card(slide, Inches(6.8), Inches(4.3), Inches(6), Inches(2.5),
         "El resultado",
         "App interactiva en Streamlit\n"
         "donde puedes meter datos\n"
         "y obtener una predicción", BLUE_LIGHT)

# ═══════════════════════════════════════════════
# SLIDE 3 — EXPLORANDO LOS DATOS (EDA)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "🔍  Explorando los datos (EDA)")
add_footer(slide)

add_image_safe(slide, f"{FIGS}/saleprice_hist.png", Inches(0.4), Inches(1.3), height=Inches(2.7))
add_image_safe(slide, f"{FIGS}/corr_heatmap.png", Inches(4.8), Inches(1.3), height=Inches(2.7))
add_image_safe(slide, f"{FIGS}/outliers_boxplot.png", Inches(9.2), Inches(1.3), height=Inches(2.7))

add_text_box(slide, Inches(0.4), Inches(4.1), Inches(4.0), Inches(0.3),
             "💰  Precios de las casas", font_size=16, color=BLUE_DARK, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.4), Inches(4.4), Inches(4.0), Inches(0.8),
             "Mayoría entre 100k y 200k\nDistribución asimétrica\n(la mayoría son casas asequibles)",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(4.8), Inches(4.1), Inches(4.0), Inches(0.3),
             "🔗  ¿Qué afecta al precio?", font_size=16, color=BLUE_DARK, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.8), Inches(4.4), Inches(4.0), Inches(0.8),
             "OverallQual (0.79) ← calidad\nGrLivArea (0.71) ← metros\nGarageCars (0.64) ← garaje",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(9.2), Inches(4.1), Inches(4.0), Inches(0.3),
             "📦  Valores faltantes", font_size=16, color=BLUE_DARK, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.2), Inches(4.4), Inches(4.0), Inches(0.8),
             "Hay columnas con muchos\nvalores vacíos, como PoolQC\n(99.5%) o MiscFeature (96.3%)",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.3), LIGHT_BG)
add_text_box(slide, Inches(0.7), Inches(5.55), Inches(12), Inches(0.3),
             "💡  Aprendizaje:", font_size=16, color=BLUE_MID, bold=True)
add_text_box(slide, Inches(0.7), Inches(5.9), Inches(12), Inches(0.7),
             "No todas las variables son importantes. Nos podemos centrar en las que más se relacionan con el precio.",
             font_size=15, color=BLACK)

# ═══════════════════════════════════════════════
# SLIDE 4 — LIMPIEZA Y PREPARACIÓN
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "🧹  Limpieza y preparación de los datos")
add_footer(slide)

add_card(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(2.3),
         "Rellenar huecos",
         "Valores vacíos:\n• Números → mediana\n• Categorías → moda", BLUE_MID)

add_card(slide, Inches(4.6), Inches(1.4), Inches(3.8), Inches(2.3),
         "Estandarizar",
         "Poner todos los números\nen la misma escala\ncon StandardScaler", BLUE_LIGHT)

add_card(slide, Inches(8.7), Inches(1.4), Inches(4.1), Inches(2.3),
         "Convertir texto a números",
         "Categorías → OneHotEncoding\n(cada opción se vuelve\nuna columna de 0/1)", BLUE_MID)

add_rect(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.8), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.3),
             "✂️  División del dataset:", font_size=20, color=BLUE_DARK, bold=True)

add_rect(slide, Inches(0.8), Inches(4.6), Inches(2.3), Inches(1.5), RGBColor(0xDD, 0xE8, 0xF4))
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(2.3), Inches(0.3),
             "80%  Train", font_size=18, color=BLUE_DARK, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(2.3), Inches(0.7),
             "1.168 casas\npara entrenar",
             font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.5), Inches(4.6), Inches(2.3), Inches(1.5), RGBColor(0xFC, 0xE4, 0xD6))
add_text_box(slide, Inches(3.5), Inches(4.8), Inches(2.3), Inches(0.3),
             "20%  Test", font_size=18, color=RGBColor(0xCC, 0x66, 0x00), bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(5.2), Inches(2.3), Inches(0.7),
             "292 casas\npara evaluar",
             font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(6.5), Inches(4.1), Inches(6), Inches(2.5),
             "¿Por qué separar?\n\n"
             "Si evaluamos con los mismos\n"
             "datos que usamos para entrenar,\n"
             "el modelo haría \"trampa\" y no\n"
             "sabríamos si funciona de verdad\n"
             "con casos nuevos.",
             font_size=15, color=GRAY)

# ═══════════════════════════════════════════════
# SLIDE 5 — LAS 6 VARIABLES CLAVE
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "⭐  Las 6 variables más importantes")
add_footer(slide)

add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             "Nos quedamos con las que más influyen en el precio. Así el modelo es más sencillo de entender.",
             font_size=18, color=GRAY)

features = [
    ("OverallQual", "0.79", "Calidad de la casa"),
    ("GrLivArea", "0.71", "Metros cuadrados habitables"),
    ("GarageCars", "0.64", "Plazas de garaje"),
    ("TotalBsmtSF", "0.61", "Tamaño del sótano"),
    ("FullBath", "0.56", "Baños completos"),
    ("YearBuilt", "0.52", "Año de construcción"),
]

y = Inches(1.8)
for i, (feat, corr, desc) in enumerate(features):
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.75), bg)
    add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(0.75), BLUE_MID)
    add_text_box(slide, Inches(0.8), y + Inches(0.08), Inches(2.5), Inches(0.55),
                 feat, font_size=20, color=BLUE_DARK, bold=True)
    add_rect(slide, Inches(3.8), y + Inches(0.12), Inches(0.9), Inches(0.45), BLUE_LIGHT)
    add_text_box(slide, Inches(3.8), y + Inches(0.14), Inches(0.9), Inches(0.45),
                 corr, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.0), y + Inches(0.08), Inches(7), Inches(0.55),
                 desc, font_size=17, color=GRAY)
    y += Inches(0.9)

add_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.6), RGBColor(0xDD, 0xEE, 0xF8))
add_text_box(slide, Inches(0.8), Inches(6.25), Inches(12), Inches(0.5),
             "💡  Con solo 6 variables podemos predecir el precio con bastante precisión. Menos es más.",
             font_size=16, color=BLUE_MID, bold=True)

# ═══════════════════════════════════════════════
# SLIDE 6 — PRIMER MODELO (BASELINE)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "🤖  Nuestro primer modelo")
add_footer(slide)

add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             "Empezamos con un modelo sencillo: Random Forest — sin ajustar, con los valores por defecto.",
             font_size=18, color=GRAY)

add_rect(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(0.3),
             "⚙️  Configuración", font_size=20, color=BLUE_DARK, bold=True)
items = [
    "  Algoritmo: RandomForestRegressor",
    "  300 árboles de decisión",
    "  Entrenado con las 6 variables",
]
add_bullets(slide, Inches(0.8), Inches(2.3), Inches(4.8), Inches(1.5), items,
            font_size=17, spacing=Pt(8))

add_rect(slide, Inches(6.5), Inches(1.8), Inches(6.3), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(0.3),
             "📊  Resultados en test", font_size=20, color=BLUE_DARK, bold=True)

metrics = [
    ("RMSE", "$29.102", "Error típico (en $)"),
    ("MAE", "$19.067", "Error absoluto medio"),
    ("R²",  "0.8896", "Precisión (1 = perfecto)"),
]
my = Inches(2.3)
for met, val, desc in metrics:
    add_text_box(slide, Inches(6.8), my, Inches(1.7), Inches(0.3),
                 met, font_size=15, color=GRAY)
    add_text_box(slide, Inches(8.5), my, Inches(1.5), Inches(0.3),
                 val, font_size=18, color=ACCENT, bold=True)
    add_text_box(slide, Inches(10.2), my, Inches(2.3), Inches(0.3),
                 desc, font_size=13, color=GRAY)
    my += Inches(0.5)

add_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.2), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(0.3),
             "📝  ¿Qué significa R² = 0.89?", font_size=18, color=BLUE_DARK, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(12), Inches(1.5),
             "El modelo explica el 89% de la variación de los precios. Dicho de otro modo:\n"
             "de cada 10 casas, el modelo acierta la tendencia de casi 9.\n\n"
             "El error típico es de unos $29.000 — para ser nuestro primer modelo con solo 6 variables, "
             "¡no está nada mal!",
             font_size=16, color=BLACK)

# ═══════════════════════════════════════════════
# SLIDE 7 — INTENTAMOS MEJORARLO (TUNING)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "🔧  Intentando mejorar el modelo")
add_footer(slide)

add_text_box(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             "Probamos a ajustar los parámetros del Random Forest para ver si podemos mejorarlo.",
             font_size=18, color=GRAY)

add_rect(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(3.0), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(0.3),
             "🔍  Cómo lo hicimos", font_size=20, color=BLUE_DARK, bold=True)
tune_items = [
    "  RandomizedSearchCV: prueba combinaciones",
    "  aleatorias (25 de 432 posibles)",
    "  3 validaciones cruzadas cada una",
    "  → 75 entrenamientos en total",
    "  Busca: profundidad, nº árboles,",
    "  muestras mínimas, vars. por división",
]
add_bullets(slide, Inches(0.8), Inches(2.3), Inches(4.8), Inches(2.0), tune_items,
            font_size=16, spacing=Pt(6))

add_rect(slide, Inches(6.5), Inches(1.8), Inches(6.3), Inches(3.0), LIGHT_BG)
add_text_box(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(0.3),
             "📊  Resultados", font_size=20, color=BLUE_DARK, bold=True)
tune_res = [
    "  RMSE: $29.724  (antes $29.102)",
    "  MAE:  $19.331  (antes $19.067)",
    "  R²:   0.8848   (antes 0.8896)",
]
add_bullets(slide, Inches(6.8), Inches(2.4), Inches(5.5), Inches(1.5), tune_res,
            font_size=17, spacing=Pt(8))

add_rect(slide, Inches(6.8), Inches(4.1), Inches(5.5), Inches(0.5), RGBColor(0xFC, 0xE4, 0xD6))
add_text_box(slide, Inches(7.0), Inches(4.15), Inches(5), Inches(0.4),
             "⚠️  El modelo por defecto funcionó mejor",
             font_size=16, color=RGBColor(0xCC, 0x66, 0x00), bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(0.3),
             "💡  ¿Qué aprendimos?", font_size=18, color=BLUE_DARK, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(12), Inches(1.0),
             "No siempre hace falta complicarse. El modelo por defecto ya era muy bueno.\n"
             "A veces, con lo básico bien hecho, se consiguen resultados excelentes.\n"
             "Esto es algo muy típico en proyectos de iniciación al ML.",
             font_size=16, color=BLACK)

# ═══════════════════════════════════════════════
# SLIDE 8 — COMPARACIÓN VISUAL
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "⚖️  Comparativa: Baseline vs Tuned")
add_footer(slide)

rows = 4
cols = 4
table_shape = slide.shapes.add_table(rows, cols, Inches(2.5), Inches(1.8), Inches(8.3), Inches(2.8))
table = table_shape.table

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.1)
table.columns[2].width = Inches(2.1)
table.columns[3].width = Inches(2.1)

headers = ["Métrica", "Baseline", "Tuned", "¿Mejoró?"]
data_rows = [
    ["RMSE", "$29.102", "$29.724", "❌  No"],
    ["MAE",  "$19.067", "$19.331", "❌  No"],
    ["R²",   "0.8896",  "0.8848",  "❌  No"],
]

for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE_DARK

for row_idx, row_data in enumerate(data_rows, 1):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = BLACK
            p.alignment = PP_ALIGN.CENTER
        bg = LIGHT_BG if row_idx % 2 == 1 else WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

add_text_box(slide, Inches(2), Inches(4.9), Inches(9), Inches(0.5),
             "El modelo baseline con valores por defecto es suficiente y más robusto",
             font_size=18, color=BLUE_MID, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(12), Inches(0.3),
             "🎯  Conclusión", font_size=18, color=BLUE_DARK, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.8),
             "Para este proyecto, con lo básico bien hecho obtenemos un resultado sólido.\n"
             "En el futuro podríamos probar otros modelos (XGBoost, Gradient Boosting) para intentar mejorarlo.",
             font_size=16, color=BLACK)

# ═══════════════════════════════════════════════
# SLIDE 9 — APP STREAMLIT
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "🌐  La app en Streamlit")
add_footer(slide)

add_image_safe(slide, f"{FIGS}/app_screenshot.png", Inches(0.5), Inches(1.3), width=Inches(8.2))

add_rect(slide, Inches(9.0), Inches(1.3), Inches(3.8), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(9.2), Inches(1.4), Inches(3.4), Inches(0.3),
             "✨  Características", font_size=18, color=BLUE_DARK, bold=True)
app_features = [
    "  Interfaz sencilla e intuitiva",
    "  Ajusta las 6 variables clave",
    "  Predicción al instante",
    "  Sin instalar nada en tu PC",
]
add_bullets(slide, Inches(9.2), Inches(1.8), Inches(3.4), Inches(1.8), app_features,
            font_size=15, spacing=Pt(6))

add_rect(slide, Inches(9.0), Inches(4.0), Inches(3.8), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(9.2), Inches(4.1), Inches(3.4), Inches(0.3),
             "🔗  Acceso", font_size=18, color=BLUE_DARK, bold=True)
add_text_box(slide, Inches(9.2), Inches(4.5), Inches(3.4), Inches(0.8),
             "Disponible online en:\nhouse-prices-\nv3nrejcyjltytdctev9hgk\n.streamlit.app",
             font_size=14, color=GRAY)

add_rect(slide, Inches(0.5), Inches(5.0), Inches(8.2), Inches(1.8), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(5.1), Inches(7.5), Inches(0.3),
             "📱  Pruébalo tú mismo", font_size=18, color=BLUE_DARK, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(7.5), Inches(1.0),
             "Mueve los deslizadores de calidad, metros cuadrados, garaje...\n"
             "y el precio se actualiza automáticamente. Así de fácil.",
             font_size=16, color=BLACK)

# ═══════════════════════════════════════════════
# SLIDE 10 — EL PIPELINE COMPLETO
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "🏗️  El pipeline completo")
add_footer(slide)

stages = [
    ("📊", "EDA", "Entender los\ndatos"),
    ("🧹", "Limpieza", "Rellenar huecos\ny estandarizar"),
    ("✨", "Selección", "Elegir las\n6 mejores"),
    ("🤖", "Modelo", "Random Forest\nentrenado"),
    ("🌐", "App", "Predicción\ninteractiva"),
]

box_w = Inches(2.2)
box_h = Inches(2.8)
start_x = Inches(0.5)
gap = Inches(0.2)
y_top = Inches(1.4)

for i, (emoji, title, desc) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    bar = BLUE_DARK if i % 2 == 0 else BLUE_MID
    add_rect(slide, x, y_top, box_w, box_h, LIGHT_BG)
    add_rect(slide, x, y_top, box_w, Inches(0.06), bar)
    add_text_box(slide, x, y_top + Inches(0.2), box_w, Inches(0.5),
                 emoji, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y_top + Inches(0.8), box_w, Inches(0.4),
                 title, font_size=18, color=bar, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y_top + Inches(1.2), box_w, Inches(0.8),
                 desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(stages) - 1:
        ax = x + box_w + Pt(2)
        add_text_box(slide, ax, y_top + Inches(0.9), Inches(0.3), Inches(0.4),
                     "→", font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.3), LIGHT_BG)
add_text_box(slide, Inches(0.8), Inches(4.55), Inches(12), Inches(0.3),
             "🖼️  Así se ven los datos en cada etapa:", font_size=18, color=BLUE_DARK, bold=True)

add_image_safe(slide, f"{FIGS}/corr_heatmap.png", Inches(0.7), Inches(4.9), height=Inches(1.7))
add_image_safe(slide, f"{FIGS}/saleprice_hist.png", Inches(3.7), Inches(4.9), height=Inches(1.7))
add_image_safe(slide, f"{FIGS}/outliers_boxplot.png", Inches(6.7), Inches(4.9), height=Inches(1.7))
add_image_safe(slide, f"{FIGS}/app_screenshot.png", Inches(9.7), Inches(4.9), height=Inches(1.7))

# ═══════════════════════════════════════════════
# SLIDE 11 — QUÉ APRENDIMOS
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "📖  Qué aprendimos con este proyecto")
add_footer(slide)

lessons = [
    ("✅  R² de 0.89", "Con solo 6 variables nuestro modelo acierta el 89% de la variación de precios"),
    ("✅  Pipeline completo", "Aprendimos a hacer todo el proceso: desde los datos hasta una app funcionando"),
    ("✅  Un buen preprocesado marca la diferencia", "Entender, limpiar y seleccionar los datos es tan importante como el modelo"),
    ("✅  La simplicidad también funciona", "A veces el modelo inicial ya es suficientemente bueno: más complejidad no siempre significa mejores resultados"),
    ("🚀  Próximos pasos", "Probar XGBoost, Gradient Boosting, más variables, log-transform..."),
]

y = Inches(1.3)
for title, desc in lessons:
    add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.0), LIGHT_BG)
    add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(1.0), BLUE_MID)
    add_text_box(slide, Inches(0.8), y + Inches(0.1), Inches(11.5), Inches(0.35),
                 title, font_size=19, color=BLUE_DARK, bold=True)
    add_text_box(slide, Inches(0.8), y + Inches(0.5), Inches(11.5), Inches(0.35),
                 desc, font_size=15, color=GRAY)
    y += Inches(1.15)

# ═══════════════════════════════════════════════
# SLIDE 12 — DESPEDIDA
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_rect(slide, 0, Inches(0), W, Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.7),
             "🙌  ¡Gracias!", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(2.6), W, Inches(0.04), ACCENT)

links = [
    ("📂  Repositorio", "github.com/Racap/house-prices"),
    ("🌐  App Streamlit", "house-prices-v3nrejcyjltytdctev9hgk.streamlit.app"),
    ("🎥  Vídeo demostración", "youtu.be/KWWt_EOOM8A"),
    ("📄  Documentación técnica", "doc.pdf"),
]

y = Inches(3.0)
for title, url in links:
    add_text_box(slide, Inches(3), y, Inches(7.5), Inches(0.4),
                 title, font_size=20, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3), y + Inches(0.35), Inches(7.5), Inches(0.35),
                 url, font_size=16, color=RGBColor(0xBB, 0xCC, 0xDD),
                 alignment=PP_ALIGN.CENTER)
    y += Inches(0.85)

add_rect(slide, 0, Inches(6.8), W, Inches(0.06), ACCENT)
add_text_box(slide, Inches(1), Inches(6.9), Inches(11), Inches(0.5),
             "¿Preguntas?  😊",
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
output_pptx = "slides.pptx"
prs.save(output_pptx)
print(f"✅ {output_pptx} generado ({len(prs.slides)} diapositivas)")
